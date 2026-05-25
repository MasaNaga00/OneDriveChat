#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Microsoft 365 Copilot 向け ナレッジ前処理スクリプト
====================================================
ローカルフォルダ内の docx / pptx / xlsx / pdf / msg / txt を、
検索されやすい Markdown に正規化して出力する。

使い方:
    python preprocess_for_copilot.py  入力フォルダ  出力フォルダ

出力された Markdown フォルダを OneDrive(将来 SharePoint) にアップし、
Copilot に参照させる。

前提: 数十ファイル / テキストPDF中心 / スキャンPDFなし
"""

import sys
import re
import datetime
from pathlib import Path

import mammoth                 # docx -> markdown
import fitz                    # PyMuPDF: pdf
import pandas as pd            # xlsx
from pptx import Presentation  # pptx
import extract_msg            # .msg

# Copilot/Graph のインデックス上限(解析テキスト4MB)に対する安全側のしきい値。
# 1ファイルがこれを超えそうなら章/トピック単位に分割する。
SPLIT_THRESHOLD_CHARS = 1_200_000  # 約1.2M文字 (日本語で約3.5MB相当の安全圏)


# ---------------------------------------------------------------------------
# 形式ごとの抽出関数。いずれも本文(Markdown文字列)を返す。
# ---------------------------------------------------------------------------
def conv_docx(path: Path) -> str:
    with open(path, "rb") as f:
        result = mammoth.convert_to_markdown(f)
    return result.value


def conv_pptx(path: Path) -> str:
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        blocks.append(f"## スライド {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                blocks.append(shape.text_frame.text.strip())
        # 表があれば拾う
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = []
                for r in tbl.rows:
                    rows.append(" | ".join(c.text.strip() for c in r.cells))
                if rows:
                    blocks.append("\n".join(rows))
        # 発表者ノート(意外と重要な情報が入っている)
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                blocks.append(f"> ノート: {note}")
    return "\n\n".join(blocks)


def conv_xlsx(path: Path) -> str:
    xls = pd.ExcelFile(path)
    blocks = []
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        if df.empty:
            continue
        blocks.append(f"## シート: {sheet}")
        # 表として読める形(Markdownテーブル)
        blocks.append(df.to_markdown(index=False))
        # 補助として「1行=1文」の文章化も付ける。検索ヒット率が上がる。
        sentences = []
        for _, row in df.iterrows():
            parts = [f"{col}は{val}" for col, val in row.items()
                     if pd.notna(val) and str(val).strip()]
            if parts:
                sentences.append("・" + "、".join(parts) + "。")
        if sentences:
            blocks.append("### 内容(文章化)\n" + "\n".join(sentences))
    return "\n\n".join(blocks)


def conv_pdf(path: Path) -> str:
    doc = fitz.open(path)
    blocks = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            blocks.append(f"## ページ {i}\n\n{text}")
    doc.close()
    body = "\n\n".join(blocks)
    # テキストがほぼ無い = スキャンPDFの可能性。警告だけ出す(今回は対象外想定)。
    if len(body.strip()) < 30:
        body += "\n\n> ⚠ テキストがほとんど抽出できませんでした。スキャンPDFの可能性があります(要OCR)。"
    return body


def conv_msg(path: Path) -> str:
    msg = extract_msg.Message(str(path))
    header = []
    if msg.subject:
        header.append(f"件名: {msg.subject}")
    if msg.date:
        header.append(f"日時: {msg.date}")
    if msg.sender:
        header.append(f"差出人: {msg.sender}")
    if msg.to:
        header.append(f"宛先: {msg.to}")
    if msg.cc:
        header.append(f"CC: {msg.cc}")
    body = (msg.body or "").strip()
    # 添付ファイル名も検索の手がかりとして残す
    attach_names = [a.longFilename or a.shortFilename
                    for a in msg.attachments if (a.longFilename or a.shortFilename)]
    msg.close()
    out = "\n".join(header) + "\n\n---\n\n" + body
    if attach_names:
        out += "\n\n### 添付ファイル\n" + "\n".join(f"- {n}" for n in attach_names)
    return out


def conv_txt(path: Path) -> str:
    for enc in ("utf-8", "cp932", "shift_jis"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return path.read_text(errors="replace")


# 拡張子 -> 変換関数のディスパッチ表
CONVERTERS = {
    ".docx": conv_docx,
    ".pptx": conv_pptx,
    ".xlsx": conv_xlsx,
    ".pdf":  conv_pdf,
    ".msg":  conv_msg,
    ".txt":  conv_txt,
}


# ---------------------------------------------------------------------------
# メタ情報(YAMLフロントマター)の付与
# ---------------------------------------------------------------------------
def build_frontmatter(src: Path, category: str = "", keywords: str = "") -> str:
    """検索の手がかりになるメタ情報を先頭に付ける。
    category / keywords は後段で Dify に自動生成させる前提で空でも可。"""
    mtime = datetime.datetime.fromtimestamp(src.stat().st_mtime)
    lines = [
        "---",
        f"title: {src.stem}",
        f"source: {src.name}",
        f"date: {mtime:%Y-%m}",
        f"category: {category}",
        f"keywords: {keywords}",
        "---",
        "",
    ]
    return "\n".join(lines)


def safe_name(name: str) -> str:
    """出力ファイル名に使えない文字を除去"""
    return re.sub(r'[\\/:*?"<>|]', "_", name)


def source_header(src: Path, part: int = 0, total: int = 1) -> str:
    """本文先頭に置く出典表記。Copilotが根拠ファイルを拾いやすくする。"""
    label = src.name
    if total > 1:
        label += f"（{part}/{total} 分割）"
    return f"> 📄 出典ファイル: **{label}**\n\n"


def source_footer(src: Path) -> str:
    """本文末尾に置く出典表記。回答の根拠として参照されやすくする。"""
    return f"\n\n---\n\n＊この内容の出典ファイル: **{src.name}**\n"


# ---------------------------------------------------------------------------
# メイン処理
# ---------------------------------------------------------------------------
def process(in_dir: Path, out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)
    targets = [p for p in in_dir.rglob("*")
               if p.is_file() and p.suffix.lower() in CONVERTERS]

    print(f"対象ファイル: {len(targets)} 件\n")
    ok, ng = 0, 0
    for src in sorted(targets):
        try:
            body = CONVERTERS[src.suffix.lower()](src)
            fm = build_frontmatter(src)
            base = safe_name(src.stem)

            # 4MBの壁対策: しきい値を超えたら見出し(##)単位で分割
            if len(fm + body) > SPLIT_THRESHOLD_CHARS:
                chunks = split_by_heading(body, SPLIT_THRESHOLD_CHARS)
                total = len(chunks)
                for idx, chunk in enumerate(chunks, 1):
                    out = out_dir / f"{base}_part{idx}.md"
                    chunk_body = (source_header(src, idx, total)
                                  + chunk
                                  + source_footer(src))
                    out.write_text(fm + chunk_body, encoding="utf-8")
                print(f"  ✓ {src.name}  → {total}分割")
            else:
                full_body = source_header(src) + body + source_footer(src)
                out = out_dir / f"{base}.md"
                out.write_text(fm + full_body, encoding="utf-8")
                print(f"  ✓ {src.name}")
            ok += 1
        except Exception as e:
            print(f"  ✗ {src.name}  ({e})")
            ng += 1

    print(f"\n完了: 成功 {ok} 件 / 失敗 {ng} 件")
    print(f"出力先: {out_dir}")


def split_by_heading(body: str, limit: int) -> list[str]:
    """## 見出し単位でまとめつつ、limit を超えないチャンクに束ねる"""
    parts = re.split(r"(?=^## )", body, flags=re.MULTILINE)
    chunks, cur = [], ""
    for part in parts:
        if len(cur) + len(part) > limit and cur:
            chunks.append(cur)
            cur = part
        else:
            cur += part
    if cur:
        chunks.append(cur)
    return chunks


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("使い方: python preprocess_for_copilot.py 入力フォルダ 出力フォルダ")
        sys.exit(1)
    process(Path(sys.argv[1]), Path(sys.argv[2]))
